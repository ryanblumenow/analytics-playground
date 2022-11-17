import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.pyplot import figure
from datetime import datetime
from streamlit_option_menu import option_menu
from io import BytesIO
from streamlit_quill import st_quill
import os
import dataingestion

#add an import to Hydralit
from hydralit import HydraHeadApp

from scipy.optimize import curve_fit

#create a wrapper class
class forecasting(HydraHeadApp):

#wrap all your code in this method and you should be done
    def run(self):
        title = '<p style="font-family:sans-serif; color:red; font-size: 39px; text-align: center;"><b>Planning and forecasting</b></p>'

        st.markdown(title, unsafe_allow_html=True)

        df, df2, branddf = dataingestion.readdata()
        print(df.head())

            ### End

            ### Enter code to test here

        df, df2, branddf = readdata()

        st.session_state['pagechoice'] = 'test'

        with st.sidebar:
            choose = option_menu("ABI Analytics", ["Keep in mind", "Show me the steps", "Give me tips", "Give feedback"],
                                icons=['key', 'bezier2', 'joystick', 'keyboard'],
                                menu_icon="app-indicator", default_index=0,
                                styles={
                "container": {"padding": "5!important", "background-color": "#fafafa"},
                "icon": {"color": "black", "font-size": "25px"}, 
                "nav-link": {"font-size": "16px", "text-align": "left", "margin":"0px", "--hover-color": "#eee"},
                "nav-link-selected": {"background-color": "#ab0202"},
            }
            )

            print(st.session_state.pagechoice)

            if choose == "Keep in mind":
                st.write("Remember to ask good questions. That is the basis of making good decisions.")

            if choose == "Show me the steps" and st.session_state.pagechoice=="test":
                st.write("Steps you should follow:")

            if choose == "Give me tips":
                st.write("Here are some tips:")

            if choose == "Give feedback":
                st.write("Give feedback")
                with st.sidebar.form(key='columns_in_form',clear_on_submit=True): #set clear_on_submit=True so that the form will be reset/cleared once it's submitted
                    rating=st.slider("Please rate the app", min_value=1, max_value=5, value=3,help='Drag the slider to rate the app. This is a 1-5 rating scale where 5 is the highest rating')
                    text=st.text_input(label='Please leave your feedback here')
                    submitted = st.form_submit_button('Submit')
                    if submitted:
                        st.write('Thanks for your feedback!')
                        st.markdown('Your Rating:')
                        st.markdown(rating)
                        st.markdown('Your Feedback:')
                        st.markdown(text)


        # print(datetime.fromisoformat('2014-03-05').timestamp())

        x1 = np.array([0.0, 1.0, 2.0, 3.0,  4.0,  5.0])
        # x = np.arange("2020-03","2020-04",dtype='datetime64[D]')
        x2  = np.array(['2014-01-05', '2014-02-05', '2014-03-05', '2014-04-05', '2014-05-05', '2014-06-05'])
        y = np.array([0.0, 0.8, 0.9, 0.1, -0.8, -1.0])
        # y = np.array([0.0, 0.8, 0.9, 1.1, 1.3, 1.5])

        explvartype = st.selectbox("Is the explanatory variable a date field or a variable?", options=['Date', "Variable"])

        if explvartype == "Date":

            x = x2

            a = []

            for i in x:
                counter = 0
                numericaldate = datetime.fromisoformat(i).timestamp()/1000000000
                a.append(numericaldate)
                counter += 1

            newx = np.array(a)
            print("newx starts here")
            print(newx)

        elif explvartype == "Variable":

            newx = x1

        # print(type(a[0]))
        # z = np.polyfit(newx, y, 3)
        try:
            z = np.polyfit(newx, y, 3)
        except:
            try:
                z = np.polyfit(newx, y, 2)
            except:
                try:
                    z = np.polyfit(newx, y, 1)
                except:
                    st.write("No polynomial fits the parameters. Please check the data since even linear extrapolation is impossible.")

        # st.write(z)

        # create polynomial
        p = np.poly1d(z)

        # plot polynomial
        fig, ax = plt.subplots(figsize=(10, 6))
        xp = np.linspace(min(newx), max(newx), 100)
        _ = plt.plot(newx, y, '.', xp, p(xp), '-')
        plt.ylim(-2,2)
        plt.xlim(min(newx), max(newx))
        # plt.xscale
        buf = BytesIO()
        fig.savefig(buf, format="png")
        st.image(buf)
        
        # show roots
        st.write("Roots (turning points) occur at the following locations:")
        st.write(p.roots)
        st.write("Prediction 1:")
        st.write(p(datetime.fromisoformat('2014-10-05').timestamp()/1000000000))

        st.subheader("Notes on forecasting")

        # Spawn a new Quill editor
        content = st_quill(placeholder="Write your notes here")

        if content:
            st.stop()

        buildpres = st.button("Build presentation")

        if buildpres == True:

            # Building Powerpoint presentation

            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.dml.color import RGBColor
            from pptx.util import Inches, Pt
            from pptx.enum.dml import MSO_THEME_COLOR
            title='   Analytics Playground\n\
            Results from analysis'
            APlogo='./Powerpoint/APlogo.png'
            ABIlogo='./Powerpoint/ABIlogo.png'
            prs = Presentation()

            # Slide 1

            # Add colour bar

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
            )
            shape.shadow.inherit = False
            fill=shape.fill
            fill.solid()
            fill.fore_color.rgb=RGBColor(255,0,0)
            shape.text= title
            line=shape.line
            line.color.rgb=RGBColor(255,0,0)
            logo1=slide.shapes.add_picture(APlogo,Inches(13.5),Inches(6.0),height=Inches(1.08),width=Inches(1.0))
            logo2=slide.shapes.add_picture(ABIlogo,Inches(14.5),Inches(5.8),height=Inches(1.5),width=Inches(1.51))

            # Add text box for results

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
            shape.shadow.inherit = False
            fill=shape.fill
            fill.solid()
            fill.fore_color.rgb=RGBColor(255,0,0)
            shape.text= "   Results from Exploratory Data Analysis"
            line=shape.line
            line.color.rgb=RGBColor(255,0,0)
            logo1=slide.shapes.add_picture(APlogo,Inches(14.5),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
            logo2=slide.shapes.add_picture(ABIlogo,Inches(15.0),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
            left = Inches(1)
            top = Inches(2)
            width = Inches(5)
            height = Inches(5)
            text_box=slide.shapes.add_textbox(left, top, width, height)
            tb=text_box.text_frame
            tb.text = content
            prg=tb.add_paragraph()
            prg.text=" "
            prg=tb.add_paragraph()
            prg.text=''
            correlpic = slide.shapes.add_picture('correl.jpg', Inches(8), Inches(1.3), height=Inches(3.7), width=Inches(6.3))
            ppspic = slide.shapes.add_picture('pps.jpg', Inches(8), Inches(5.1), height=Inches(3.7), width=Inches(7.3))

            prs.save('EDA_presentation.pptx')

            os.startfile("EDA_presentation.pptx")

            