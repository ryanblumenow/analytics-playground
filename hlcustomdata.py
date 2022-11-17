from os import curdir
import hydralit as hy
from numpy.core.fromnumeric import var
import streamlit
import streamlit as st
import sys
# from streamlit import cli as stcli
from PIL import Image
from functions import *
import streamlit.components.v1 as components
import pandas as pd
import numpy as np
import statsmodels.api as sm
from numpy import mean
from numpy import std
from sklearn.datasets import make_classification
from sklearn.model_selection import cross_val_score
from sklearn.model_selection import RepeatedStratifiedKFold
from sklearn.linear_model import LogisticRegression
from matplotlib import pyplot
from sklearn.model_selection import train_test_split
import matplotlib.pyplot as plt
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.metrics import silhouette_score
from scipy.spatial.distance import cdist
import seaborn as sns
from io import BytesIO
from statsmodels.formula.api import ols
# from streamlit.session_state import SessionState
# import tkinter
import matplotlib
# matplotlib.use('TkAgg')
# matplotlib.use('Agg')
from sklearn.neural_network import MLPClassifier
from sklearn.metrics import accuracy_score
from mpl_toolkits.mplot3d import Axes3D
from matplotlib import cm
from matplotlib.ticker import LinearLocator, FormatStrFormatter
from sklearn.tree import DecisionTreeRegressor, plot_tree
import sklearn
from sklearn.datasets import make_classification
from sklearn.ensemble import RandomForestClassifier
from sklearn.datasets import make_regression
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import RepeatedKFold
import time
from scipy.stats import pearsonr
from scipy.stats import spearmanr
import dtale
from dtale.views import startup
from dtale.app import get_instance
import webbrowser
import dtale.global_state as global_state
import dtale.app as dtale_app
from matplotlib.pyplot import hist
from scipy import stats as stats
from bioinfokit.analys import stats
from streamlit_option_menu import option_menu
# import pandas_profiling
from streamlit_pandas_profiling import st_profile_report
from sklearn.preprocessing import LabelEncoder
from streamlit_quill import st_quill
import dataingestion
import ppscore as pps
from streamlit_option_menu import option_menu

#add an import to Hydralit
from hydralit import HydraHeadApp

#create a wrapper class
class customdataenv(HydraHeadApp):

#wrap all your code in this method and you should be done
    def run(self):
        # sourcery skip: remove-redundant-fstring, replace-interpolation-with-fstring
        #-------------------existing untouched code------------------------------------------

        # st.title('Cheat sheets and guides')

        title = '<p style="font-family:sans-serif; color:red; font-size: 39px; text-align: center;"><b>Custom data environment</b></p>'
        st.markdown(title, unsafe_allow_html=True)

        # df, df2, branddf = dataingestion.readdata()
        # print(df.head())

        ### Enter code to test here

        st.session_state['pagechoice'] = 'test'

        with st.sidebar:
            choose = option_menu("ABI Analytics", ["Keep in mind", "Show me the steps", "Give me tips", "Give feedback"],
                                icons=['key', 'bezier2', 'joystick', 'keyboard'],
                                menu_icon="app-indicator", default_index=0,
                                styles={
                "container": {"padding": "5!important", "background-color": "#fafafa"},
                "icon": {"color": "black", "font-size": "25px"}, 
                "nav-link": {"font-size": "16px", "text-align": "left", "margin":"0px", "--hover-color": "#eee"},
                "nav-link-selected": {"background-color": "#ab0202"},
            }
            )

            print(st.session_state.pagechoice)

            if choose == "Keep in mind":
                st.write("Remember to ask good questions. That is the basis of making good decisions.")

            if choose == "Show me the steps" and st.session_state.pagechoice=="test":
                st.write("Steps you should follow:")

            if choose == "Give me tips":
                st.write("Here are some tips:")

            if choose == "Give feedback":
                st.write("Give feedback")
                with st.sidebar.form(key='columns_in_form',clear_on_submit=True): #set clear_on_submit=True so that the form will be reset/cleared once it's submitted
                    rating=st.slider("Please rate the app", min_value=1, max_value=5, value=3,help='Drag the slider to rate the app. This is a 1-5 rating scale where 5 is the highest rating')
                    text=st.text_input(label='Please leave your feedback here')
                    submitted = st.form_submit_button('Submit')
                    if submitted:
                        st.write('Thanks for your feedback!')
                        st.markdown('Your Rating:')
                        st.markdown(rating)
                        st.markdown('Your Feedback:')
                        st.markdown(text)

        uploaded_file = st.file_uploader("Choose a file")

        if uploaded_file is not None:

            if 'load_csv' in st.session_state:
                dataframe=st.session_state.load_csv
                st.write(uploaded_file.name + " " +  "is loaded") 

            else:
                dataframe = pd.read_csv(uploaded_file, dtype='unicode')
                st.session_state.load_csv = dataframe

            # dataframe = pd.read_csv(uploaded_file, dtype='unicode')

            with st.spinner("Auto-encoding categorical variables for analysis"):

                #Auto encodes any dataframe column of type category or object.
                # def dummyEncode(df):
                columnsToEncode = list(dataframe.select_dtypes(include=['category','object']))
                le = LabelEncoder()
                for feature in columnsToEncode:
                    try:
                        dataframe[feature] = le.fit_transform(dataframe[feature])
                    except:
                        print('Error encoding '+feature)
                # print(dataframe)
                dataframe.to_csv('dataframe.csv', index=False)

            edaprofiling = st.expander("Profile of dataset " + uploaded_file.name, expanded=False)

            with edaprofiling:

                # @st.cache(allow_output_mutation=True)
                # def gen_profile_report(dataframe, *report_args, **report_kwargs):
                #     return dataframe.profile_report(*report_args, **report_kwargs)

                # pr = gen_profile_report(dataframe, explorative=True)

                # st_profile_report(pr)

                dataframe2 = dataframe.sample(1000).reset_index(drop=True)
                st.write("Shape of dataset:")
                st.write(dataframe2.shape)

                from pandas_profiling import ProfileReport

                @st.cache(allow_output_mutation=True)
                def gen_profile_report(dataframe2, *report_args, **report_kwargs):
                    # return dataframe2.profile_report(*report_args, **report_kwargs)
                    return ProfileReport(dataframe2, *report_args, **report_kwargs)
                    # from pandas_profiling import ProfileReport
                    # return ProfileReport(dataframe2, *report_args, **report_kwargs)                    

                pr = gen_profile_report(dataframe2, minimal=True, title="Data profile",
                dataset={
                "description": "This profiling report shows an overview of the data",
                "copyright_holder": "Analytics Playground",
                "copyright_year": "2022",
                "url": "https://www.ryanblumenow.com"},
                vars={"num": {"low_categorical_threshold": 0}},
                plot={'histogram':{'bayesian_blocks_bins': False}})

                st_profile_report(pr)

                st.write("")

                dataframe2 = dataframe2.dropna()

                st.subheader("Correlation between variables in dataset")

                st.write(dataframe2.corr(method = 'pearson'))

                # corrmap = sns.heatmap(corr, vmin=0, vmax=1, cmap="Blues", linewidths=0.5, annot=True)
                plt.figure(figsize=(16, 6))
                corrheatmap = sns.heatmap(dataframe2.corr(), vmin=-1, vmax=1, annot=False, cmap='BrBG')
                corrheatmap.set_title('Correlation heatmap', fontdict={'fontsize':12}, pad=12)
                
                plt.savefig('customcorr.png', dpi=300, bbox_inches='tight')

                st.image('./customcorr.png')

                st.write("")

                # corrmap = sns.heatmap(dataframe2.corr(),annot=True)

                st.write("")

                st.subheader("PPS scores for variables in dataset")

                # matrix_pps = pps.matrix(dataframe2)

                matrix_pps = pps.matrix(dataframe2)[['x', 'y', 'ppscore']].pivot(columns='x', index='y', values='ppscore')

                st.write(matrix_pps)

                plt.figure(figsize=(16, 6))
                sns.heatmap(matrix_pps, vmin=0, vmax=1, cmap="Blues", linewidths=0.5, annot=False)

                # ppsfigure = ppsmap.get_figure()    
                # ppsfigure.savefig('custompps.jpg', dpi=400)

                plt.savefig('custompps.png', dpi=300, bbox_inches='tight')

                st.image('./custompps.jpg')

            dtaleexpander = st.expander("Dataset analysis", expanded=False)

            with dtaleexpander:

                startup(data_id="1", data=dataframe)

                if get_instance("1") is None:
                    startup(data_id="1", data=dataframe)

                d=get_instance("1")

                # webbrowser.open_new_tab('http://localhost:8501/dtale/main/1') # New window/tab
                # components.html("<iframe src='/dtale/main/1' />", width=1000, height=300, scrolling=True) # Element
                html = f"""<iframe src="/dtale/main/1" height="1000" width="1400"></iframe>""" # Iframe
                # html = "<a href='/dtale/main/1' target='_blank'>Dataframe 1</a>" # New tab link

                st.markdown(html, unsafe_allow_html=True)

                checkbtn = st.button("Validate data")

                if checkbtn == True:
                    df_amended = get_instance(data_id="1").data # The amended dataframe ready for upsert
                    st.write("Sample of amended data:")
                    st.write("")
                    st.write(dataframe_amended.head(5))

                clearchanges = st.button("Clear changes made to data")
                if clearchanges == True:
                    global_state.cleanup()

            st.subheader("Notes on custom data analysis")

            # Spawn a new Quill editor
            customcontent = st_quill(placeholder="Write your notes here")

            st.session_state.customdatanotes = customcontent

            if customcontent:
                st.stop()

            st.write("Take a screenshot of the correlation heatmap and save it as 'customcorrel.jpg', do the same for the PPS heatmap and call it 'custompps.jpg'.")

            buildpres = st.button("Build presentation")

            if buildpres == True:

                # Building Powerpoint presentation

                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.dml.color import RGBColor
                from pptx.util import Inches, Pt
                from pptx.enum.dml import MSO_THEME_COLOR
                title='   Analytics Playground\n\
                Results from analysis'
                APlogo='./Powerpoint/APlogo.png'
                ABIlogo='./Powerpoint/ABIlogo.png'
                prs = Presentation()

                # Slide 1

                # Add colour bar

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
                )
                shape.shadow.inherit = False
                fill=shape.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor(255,0,0)
                shape.text= title
                line=shape.line
                line.color.rgb=RGBColor(255,0,0)
                logo1=slide.shapes.add_picture(APlogo,Inches(13.5),Inches(6.0),height=Inches(1.08),width=Inches(1.0))
                logo2=slide.shapes.add_picture(ABIlogo,Inches(14.5),Inches(5.8),height=Inches(1.5),width=Inches(1.51))

                # Slide 2

                # Add text box for results

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
                shape.shadow.inherit = False
                fill=shape.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor(255,0,0)
                shape.text= "   Results from custom data analysis"
                line=shape.line
                line.color.rgb=RGBColor(255,0,0)
                logo1=slide.shapes.add_picture(APlogo,Inches(14.5),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                logo2=slide.shapes.add_picture(ABIlogo,Inches(15.0),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                left = Inches(1)
                top = Inches(2)
                width = Inches(5)
                height = Inches(5)
                text_box=slide.shapes.add_textbox(left, top, width, height)
                tb=text_box.text_frame
                tb.text = st.session_state.customdatanotes
                prg=tb.add_paragraph()
                prg.text=" "
                prg=tb.add_paragraph()
                prg.text=''
                correlpic = slide.shapes.add_picture('customcorrel.jpg', Inches(8), Inches(1.3), height=Inches(3.7), width=Inches(6.3))
                ppspic = slide.shapes.add_picture('custompps.jpg', Inches(8), Inches(5.1), height=Inches(3.7), width=Inches(7.3))

                prs.save('Custom_data_presentation.pptx')

                os.startfile("Custom_data_presentation.pptx")

