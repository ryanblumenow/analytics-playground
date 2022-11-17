import streamlit as st
from streamlit_option_menu import option_menu
import base64
from email import header
from html.entities import html5
# from importlib.resources import read_binary
import hydralit as hy
# from markdown import markdown
from numpy.core.fromnumeric import var
import streamlit
import streamlit as st
import sys
#from streamlit import cli as stcli
from PIL import Image
from functions import *
import streamlit.components.v1 as components
import pandas as pd
import numpy as np
import statsmodels.api as sm
from numpy import mean
from numpy import std
from sklearn.datasets import make_classification
from sklearn.model_selection import cross_val_score
from sklearn.model_selection import RepeatedStratifiedKFold
from sklearn.linear_model import LogisticRegression
from matplotlib import pyplot
from sklearn.model_selection import train_test_split
import matplotlib.pyplot as plt
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.metrics import silhouette_score
from scipy.spatial.distance import cdist
import seaborn as sns
from io import BytesIO
from statsmodels.formula.api import ols
# from streamlit.session_state import SessionState
# import tkinter
import matplotlib
# matplotlib.use('TkAgg')
# matplotlib.use('Agg')
from sklearn.neural_network import MLPClassifier
from sklearn.metrics import accuracy_score
from mpl_toolkits.mplot3d import Axes3D
from matplotlib import cm
from matplotlib.ticker import LinearLocator, FormatStrFormatter
from sklearn.tree import DecisionTreeRegressor, plot_tree
import sklearn
from sklearn.datasets import make_classification
from sklearn.ensemble import RandomForestClassifier
from sklearn.datasets import make_regression
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import RepeatedKFold
import time
from scipy.stats import pearsonr
from scipy.stats import spearmanr
import dtale
from dtale.views import startup
from dtale.app import get_instance
import webbrowser
import dtale.global_state as global_state
import dtale.app as dtale_app
from matplotlib.pyplot import axis, hist
from scipy import stats as stats
from bioinfokit.analys import stat
from statsmodels.stats.anova import AnovaRM
import statsmodels.api as sm
from statsmodels.graphics.factorplots import interaction_plot
from sklearn.decomposition import PCA
from st_click_detector import click_detector
from streamlit.components.v1 import html
import streamlit.components.v1 as components
import pandas as pd
import streamlit as st
import click_image_copy_from_demo
from dtale.views import startup
from streamlit_quill import st_quill
# import pandas_profiling
from streamlit_pandas_profiling import st_profile_report
import dataingestion
from streamlit_option_menu import option_menu

#add an import to Hydralit
from hydralit import HydraHeadApp

#create a wrapper class
class autoanalytics(HydraHeadApp):

#wrap all your code in this method and you should be done
    def run(self):
        #-------------------existing untouched code------------------------------------------

        title = '<p style="font-family:sans-serif; color:red; font-size: 39px; text-align: center;"><b>Automated analytics flow</b></p>'
        st.markdown(title, unsafe_allow_html=True)

        if 'edanotes' not in st.session_state:
            st.session_state.edanotes = "Write your notes on EDA here"
        
        if 'correlnotes' not in st.session_state:
            st.session_state.correlnotes = "Write your notes on correlation analysis here"

        if 'hyptestingnotes' not in st.session_state:
            st.session_state.hyptestingnotes = "Write your notes on hypothesis testing here"

        if 'dimrednotes' not in st.session_state:
            st.session_state.dimrednotes = "Write your notes on dimension reduction here"

        if 'regnotes' not in st.session_state:
            st.session_state.regnotes = "Write your notes on regression analysis here"

        if 'clusternotes' not in st.session_state:
            st.session_state.clusternotes = "Write your notes on cluster analysis here"

        if 'conjointnotes' not in st.session_state:
            st.session_state.conjointnotes = "Write your notes on conjoint analysis here"

        if 'dectreenotes' not in st.session_state:
            st.session_state.dectreenotes = "Write your notes on decision tree analysis here"

        ### From Jupyter - 0. Prepare the data

        df, df2, branddf = dataingestion.readdata()
        print(df.head())
        
        vars = list(df2.columns.values.tolist())

        st.session_state['pagechoice'] = 'auto analytics'

        with st.sidebar:
            choose = option_menu("ABI Analytics", ["Keep in mind", "Show me the steps", "Give me tips", "Give feedback"],
                                icons=['key', 'bezier2', 'joystick', 'keyboard'],
                                menu_icon="app-indicator", default_index=0,
                                styles={
                "container": {"padding": "5!important", "background-color": "#fafafa"},
                "icon": {"color": "black", "font-size": "25px"}, 
                "nav-link": {"font-size": "16px", "text-align": "left", "margin":"0px", "--hover-color": "#eee"},
                "nav-link-selected": {"background-color": "#ab0202"},
            }
            )

            print(st.session_state.pagechoice)
            
            if choose == "Keep in mind":
                st.write("Remember to ask good questions. That is the basis of making good decisions.")

            if choose == "Show me the steps" and st.session_state.pagechoice=="test":
                st.write("Steps you should follow:")

            if choose == "Give me tips":
                st.write("Here are some tips:")

            if choose == "Give feedback":
                st.write("Give feedback")
                with st.sidebar.form(key='columns_in_form',clear_on_submit=True): #set clear_on_submit=True so that the form will be reset/cleared once it's submitted
                    rating=st.slider("Please rate the app", min_value=1, max_value=5, value=3,help='Drag the slider to rate the app. This is a 1-5 rating scale where 5 is the highest rating')
                    text=st.text_input(label='Please leave your feedback here')
                    submitted = st.form_submit_button('Submit')
                    if submitted:
                        st.write('Thanks for your feedback!')
                        st.markdown('Your Rating:')
                        st.markdown(rating)
                        st.markdown('Your Feedback:')
                        st.markdown(text)


        value = click_image_copy_from_demo.st_click_image()
        if value is None:
            # st.stop()
            pass

        st.success("{} selected".format(value))

        col1, col2, col3 = st.columns([1,0.6,1])

        clicked = ""

        with col1:
            pass

        with col2:

            content2 = """
            <link rel="stylesheet" href="https://www.w3schools.com/w3css/4/w3.css">
                <style>
                .zoom {
                padding: 10px;
                width: 290px;
                height: 180px;
                transition: transform .21s; /* Animation */
                margin: 0 auto;
                }
                .zoom:hover {
                transform: scale(1.1); /* (110% zoom - Note: if the zoom is too large, it will go outside of the viewport) */
                }
                </style>
                <div class="w3-container">
                <a href='#' id='Get started'><img width='21%' src='https://i.postimg.cc/2yyc69HY/analyticslbl.jpg' class='zoom w3-round-xxlarge w3-hover-opacity'></a>
                <!--<div class="w3-display-bottom w3-large"><b><h3>Click to start</h3></b></div>-->
                </div>
                """

            clicked = click_detector(content2)

            st.markdown(f"**{clicked} selected - begin analysis below.**" if clicked != "" else "         **Select a methodology and click 'Get started'!**")

        with col3:
            pass

        if clicked == "Get started":

            if value=='EDA':
                
                st.subheader("Exploratory data analysis")

                with st.spinner("Analyzing and summarizing dataset and generating dataset profile"):

                    my_bar = st.progress(0)

                    for percent_complete in range(100):
                        time.sleep(0.01)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    edaenv = st.expander("Guidance on EDA", expanded=False)

                    with edaenv:

                        st.info("User guide")

                        def show_pdf(file_path):
                            # Opening tutorial from file path
                            with open(file_path, "rb") as f:
                                base64_pdf = base64.b64encode(f.read()).decode('utf-8')

                            # Embedding PDF in HTML
                            pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="1500" height="800" type="application/pdf"></iframe>'

                            # Displaying File
                            st.markdown(pdf_display, unsafe_allow_html=True)

                        col1, col2,col3= st.columns(3)
                        with col1:  
                            if st.button('Read PDF tutorial',key='1'):            
                                show_pdf('.\Automated flow\DtaleInstructions-compressed.pdf')
                        with col2:
                            st.button('Close PDF tutorial',key='2')                   
                        with col3:
                            with open(".\Automated flow\DtaleInstructions-compressed.pdf", "rb") as pdf_file:
                                PDFbyte = pdf_file.read()
                            st.download_button(label="Download PDF tutorial", key='3',
                                    data=PDFbyte,
                                    file_name="EDA Instructions.pdf",
                                    mime='application/octet-stream')

                    datadescrip = st.expander("Description of data")

                    with datadescrip:

                        st.write(df.describe(include='all'))
                        
                    edaprofiling = st.expander("Profile of dataset", expanded=False)
                    
                    with edaprofiling:
                    
                        # @st.cache(allow_output_mutation=True)
                        # def gen_profile_report(df, *report_args, **report_kwargs):
                        #     return df.profile_report(*report_args, **report_kwargs)

                        # pr = gen_profile_report(df, explorative=True)

                        # st_profile_report(pr)

                        @st.cache(allow_output_mutation=True)
                        def gen_profile_report(df, *report_args, **report_kwargs):
                            return df.profile_report(*report_args, **report_kwargs)

                        pr = gen_profile_report(df, explorative=True, title="Data profile",
                        dataset={
                        "description": "This profiling report shows an overview of the data",
                        "copyright_holder": "Analytics Playground",
                        "copyright_year": "2022",
                        "url": "https://www.ryanblumenow.com"}, vars={"num": {"low_categorical_threshold": 0}} )

                        st_profile_report(pr)

                startup(data_id="1", data=df2.sample(15000)) # All records, no OHE

                if get_instance("1") is None:
                    startup(data_id="1", data=df.sample(15000))

                d=get_instance("1")

                # webbrowser.open_new_tab('http://localhost:8501/dtale/main/1') # New window/tab
                # components.html("<iframe src='/dtale/main/1' />", width=1000, height=300, scrolling=True) # Element
                html = f"""<iframe src="/dtale/main/1" height="1000" width="1400"></iframe>""" # Iframe
                # html = "<a href='/dtale/main/1' target='_blank'>Dataframe 1</a>" # New tab link

                st.markdown(html, unsafe_allow_html=True)

                # d = dtale.show(pd.DataFrame(df2.sample(1000)))
                st.session_state.corr_img = d.get_corr_matrix()
                st.session_state.corr_df = d.get_corr_matrix(as_df=True)
                st.session_state.pps_img = d.get_pps_matrix()
                st.session_state.pps_df = d.get_pps_matrix(as_df=True)

                print(st.session_state.corr_df)

                checkbtn = st.button("Validate data")

                if checkbtn == True:
                    df_amended = get_instance(data_id="1").data # The amended dataframe ready for upsert
                    st.write("Sample of amended data:")
                    st.write("")
                    st.write(df_amended.head(5))

                clearchanges = st.button("Clear changes made to data")
                if clearchanges == True:
                    global_state.cleanup()

                st.write("")
                
                st.subheader("Notes on EDA")

                # Spawn a new Quill editor
                st.subheader("Notes on exploratory data analysis")
                edacontent = st_quill(placeholder="Write your notes here", value=st.session_state.edanotes, key="edaquill")

                st.session_state.edanotes = edacontent

                st.write("Exploratory data analysis took ", time.time() - start_time, "seconds to run")

            elif value == "Correlation analysis":

                st.subheader("Correlation analysis")

                with st.spinner('Please wait while we conduct the correlation analysis'):

                    my_bar = st.progress(0)

                    time.sleep(10)

                    for percent_complete in range(100):
                        time.sleep(0.1)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    x=0
                    y=0

                    x = st.selectbox("Please choose first variable:", vars, key="correlx", index=x)
                    y = st.selectbox("Please choose second variable:", vars, key="correly", index=y)

                    # var1 = df[x].iloc[1:101]
                    # var2 = df[y].iloc[1:101]

                    var1 = df[x].sample(10000)
                    var2 = df[y].sample(10000)

                    st.subheader("Correlation between chosen 2 variables:" + " " + x + " and " + y)

                    st.write("Correlation coefficient: ", var1.corr(var2))
                    st.write("Pearson correlation coefficient: ", pearsonr(var1, var2))
                    st.write("Spearman correlation coefficient: ", spearmanr(var1, var2))

                    ### From Jupyter: 1. Hypothesis testing: correlation

                    st.subheader("Correlation matrix, Pearson")

                    st.write(df.corr(method = 'pearson'))

                    with st.expander("Interpretation guide"):
                        st.write("Variables within a dataset can be related for many reasons.\n\n"
                        "For example:\n"
                        "One variable could cause or depend on the values of another variable.\n"
                        "One variable could be lightly associated with another variable.\n"
                        "Two variables could depend on a third unknown variable.\n\n"
                        "It can be useful in data analysis and modeling to better understand the relationships between variables. The statistical relationship between two variables is referred to as their correlation.\n"
                        "A correlation could be positive, meaning both variables move in the same direction, or negative, meaning that when one variable's value increases, the other variables' values decrease. Correlation can also be neutral or zero, meaning that the variables are unrelated.")
                        st.info("Positive Correlation: both variables change in the same direction.\n\n"
                        "Neutral Correlation: No relationship in the change of the variables.\n\n"
                        "Negative Correlation: variables change in opposite directions.\n\n"
                        "A correlation coefficient closer to 1 or to -1 indicates stronger relationships between the variables, while a correlation coefficient closer to 0 indicates a lesser relationship between the variables.")
                        st.write("The performance of some algorithms can deteriorate if two or more variables are tightly related, called multicollinearity. An example is linear regression, where one of the offending correlated variables should be removed in order to improve the skill of the model.\n\n"
                        "We may also be interested in the correlation between input variables with the output variable in order provide insight into which variables may or may not be relevant as input for developing a model.\n"
                        "The structure of the relationship may be known, e.g. it may be linear, or we may have no idea whether a relationship exists between two variables or what structure it may take. Depending what is known about the relationship and the distribution of the variables, different correlation scores can be calculated.")

                    ### End

                    st.write("")

                    # Spawn a new Quill editor
                    st.subheader("Notes on correlation analysis")
                    correlcontent = st_quill(placeholder="Write your notes here", value = st.session_state.correlnotes, key="correlquill")

                    st.session_state.correlnotes = correlcontent

                    st.write("Correlation analysis took ", time.time() - start_time, "seconds to run")

            elif value == "Hypothesis testing":

               st.subheader("Hypothesis testing on characteristics with optional disaggregation by brand or brand groups")
               
               with st.spinner('Please wait while we conduct hypothesis testing'):

                    my_bar = st.progress(0)

                    for percent_complete in range(100):
                        time.sleep(0.1)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    st.write("**Background to the model**")

                    col1, col2, col3 = st.columns(3)
                    with col1:
                        pass
                    with col2:
                        st.image("ttest.png")
                    with col3:
                        pass

                    with st.expander("Assumptions of the one-sample t-test"):
                        st.success("One-sample t-test")
                        st.write("Student's t-test or t-test is a parametric inferential statistical method used for comparing the means between two different groups (two-sample t-test) or with the specific value (one-sample t-test).")
                        st.write("This test works well for hypothesis testing, particularly amongst data with low sample sizes (in contrast to the z-test).")
                        st.write("One Sample t-test (single sample t-test) is used to compare the sample mean (a random sample from a population) with the specific value (hypothesized or known mean of the population).")
                        st.info("Dependent variable should have an approximately normal distribution (Shapiro-Wilks Test).")
                        st.info("Observations are independent of each other.")
                        st.info("Note: One sample t-test is relatively robust to the assumption of normality when the sample size is large (n ≥ 30).")
                        st.warning("Null hypothesis: Sample mean is equal to the hypothesized or known population mean.")
                        st.warning("Alternative hypothesis: Sample mean is not equal to the hypothesized or known population mean (two-tailed or two-sided).")
                        st.warning("Alternative hypothesis: Sample mean is either greater or lesser to the hypothesized or known population mean (one-tailed or one-sided).")

                    with st.expander("Assumptions of the two-sample t-test"):
                        st.success("One-sample t-test")
                        st.write("The two-sample (unpaired or independent) t-test compares the means of two independent groups, determining whether they are equal or significantly different. In two sample t-test, usually, we compute the sample means from two groups and derives the conclusion for the population’s means (unknown means) from which two groups are drawn.")
                        st.info("Observations in two groups have an approximately normal distribution (Shapiro-Wilks Test).")
                        st.info("Homogeneity of variances (variances are equal between treatment groups) (Levene or Bartlett Test).")
                        st.info("The two groups are sampled independently from each other from the same population.")
                        st.info("Note: Two sample t-test is relatively robust to the assumption of normality and homogeneity of variances when sample size is large (n ≥ 30) and there are equal number of samples (n1 = n2) in both groups.")
                        st.warning("Null hypothesis: Two group means are equal.")
                        st.warning("Alternative hypothesis: Two group means are different (two-tailed or two-sided).")
                        st.warning("Alternative hypothesis: Mean of one group either greater or lesser than another group (one-tailed or one-sided).")

                    with st.expander("Sample size considerations for t-test"):
                        st.write("The t-test can be applied for the extremely small sample size (n = 2 to 5) provided the effect size is large and data follows the t-test assumptions. Remember, a larger sample size is preferred over small sample sizes.")
                        st.write("t-test is relatively robust to the assumption of normality and homogeneity of variances when the sample size is large (n ≥ 30).")

                    brands = pd.Series(df2['BRAND'].drop_duplicates()).sort_values().tolist() # Faster than .unique()
                    brandsmod = ["None"] + brands
                    
                    st.markdown("**Please choose brands. If only one brand is chosen, a one sample t-test will be performed. If two brands are chosen, a two sample t-test will be performed.**")

                    brand1 = st.selectbox("Please select first brand of interest.", brandsmod, key="brand1")
                    brand2 = st.selectbox("Please select second brand of interest.", brandsmod, key="brand2")

                    ttestvars = list(df.drop(['BRAND', 'BRANDNAME'], axis=1).columns)
                    # ttestvars2 = ["NONE"] + ttestvars
                    ttestvars.insert(0, "NONE")

                    var1 = st.selectbox("Please select variable for analysis", ttestvars)
                    # var2 = st.selectbox("Please select second variable for analysis", ttestvars2)

                    if var1 != "NONE":

                        mask1 = df.BRANDNAME==brand1
                        dfa = df[mask1]

                        mask2 = df.BRANDNAME==brand2
                        dfb = df[mask2]

                        result1s, result2s = [0,0]

                        from scipy import stats

                        if brand1 != "None":

                            if brand2 == "None":

                                # 1-sample t-test using scipy
                                a =  dfa[var1].to_numpy()
                                print(a)
                                popmean = df[var1].mean()
                                # use parameter "alternative" for two-sided or one-sided test
                                st.write("**Results of one sample t-test:**")
                                result1s = stats.ttest_1samp(a=a, popmean=popmean)
                                st.write(result1s)

                                res = stat()
                                res.ttest(df=dfa, test_type=1, res=var1, mu=5)
                                st.text(res.summary)

                            else:

                                # 2-sample t-test

                                options = [brand1, brand2]

                                # dfc = df.loc[(df["BRANDNAME"] == brand1) & (df["BRANDNAME"] == brand2)]

                                dfc = df[df['BRANDNAME'].isin(options)]

                                # dfc = df.query("BRAND==" + brand1 + " & BRAND==" + brand2)

                                a = dfa[var1].to_numpy()
                                b = dfb[var1].to_numpy()

                                result2s = stats.ttest_ind(a=a, b=b, equal_var=False)

                                st.write(result2s)

                                dfc.loc[dfc["BRANDNAME"] == brand1, "BRAND"] = 1
                                dfc.loc[dfc["BRANDNAME"] == brand2, "BRAND"] = 2

                                res2=stat()
                                res2.ttest(df=dfc, xfac="BRANDNAME", res=var1, evar=False, test_type=2) # evar=False for unequal variance t-test (Welch)
                                st.text(res2.summary)

                                st.write("**Description of chosen variable for each brand**")

                                # st.write(df[var1].loc[df["BRANDNAME" == brand1]].describe())
                                st.write(dfc[[var1, "BRAND"]].groupby("BRAND").describe())

                                print(brand1, brand2)

                                # cond = {'a': brand1, 'b': brand2}
                                # hist = df.loc[(df.BRANDNAME == cond['a']) & (df.BRANDNAME == cond['b']), ['x','y']].plot(title='a: {a}, b: {b} distribution histogram'.format(**cond))

                                # fig, ax = df.plot.hist(bins=12, alpha=0.5)

                                subheader = ("**Histogram of {} for each brand chosen**").format(var1)

                                st.write(subheader)

                                hist = dfc.pivot(columns='BRANDNAME').AGE.plot(kind = 'hist', stacked=True, alpha=0.3, title="Frequency distribution histogram for {}".format(var1)).figure

                                buf = BytesIO()
                                hist.savefig(buf, format="png")
                                # st.pyplot(fig)
                                st.image(buf)

                            with st.expander("Interpretation guide"):
                                st.info("The higher the t-score, and the lower the p-value, the more able you are to reject the null hypothesis in favour of the alternative hypothesis.")
                                st.info("Couple this insight with a view of the underlying distributions, to see how the hypotheses function.")
                                # print(result1s, result2s, result2s[1])
                                if brand2 == "None":
                                    res1 = result1s.pvalue
                                    if res1 <= 0.1:
                                        st.success(f"Our p-value is less than 0.1, so we can reject the null hypothesis (10% level of significance) and conclude that our consumers for {brand1} are different to the total population of consumers across different brands with respect to {var1}. Consult the histogram to see how.")
                                    if res1 > 0.1:
                                        st.error(f"Our p-value is more than 0.1, so we can reject the null hypothesis (10% level of significance) and conclude that our consumers for {brand1} are similar to the total population of consumers across different brands with respect to {var1}.")
                                else:
                                    res2 = result2s.pvalue
                                    if res2 <= 0.1:
                                        st.success(f"Our p-value is less than 0.1, so we can reject the null hypothesis (10% level of significance) and conclude that our consumers for each chosen brand are different with respect to {var1}. Consult the histogram to see how.")
                                    if res2 > 0.1:
                                        st.error(f"Our p-value is more than 0.1, so we can reject the null hypothesis (10% level of significance) and conclude that our consumers for each of the chosen brands are similar with respect to {var1}.")

                    else:
                        pass

                    # Spawn a new Quill editor
                    st.subheader("Notes on hypothesis analysis")
                    hyptestingcontent = st_quill(placeholder="Write your notes here", key="hyptestingquill")

                    st.session_state.hyptestingnotes = hyptestingcontent

                    st.write("Hypothesis testing took ", time.time() - start_time, "seconds to run")

            elif value == "Dimension reduction":

                st.subheader("Dimension reduction")

                with st.spinner('Please wait while we set up the dimension reduction'):

                    my_bar = st.progress(0)

                    time.sleep(5)

                    for percent_complete in range(100):
                        time.sleep(0.1)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    st.subheader("Dimension reduction using random forest analysis")

                    dimredvars = ["None"] + list(df.columns)

                    dimredvar = st.selectbox("Please choose variable of interest:", dimredvars, key="dimredy")

                    if dimredvar != "None":

                        # y = dfdimred.loc[:,dimredvar]
                        # y = df.loc[0:10,["BRAND"]]

                        y = tuple(list(df[dimredvar])) # Make hashable

                        X = df.drop([dimredvar, 'BRANDNAME'], axis=1)

                        rf = RandomForestClassifier(n_estimators=100, max_depth=3,
                                            bootstrap=True, n_jobs=-1,
                                            random_state=0)
                        rf.fit(X, y)

                        feature_imp = pd.Series(rf.feature_importances_, 
                                                index=X.columns).sort_values(ascending=False)

                        st.write('Feature importances: ', rf.feature_importances_)

                        col1, col2, col3 = st.columns([0.5, 3, 0.5])
                        with col1:
                            pass
                        with col2:
                            # SMALL_SIZE = 5
                            # MEDIUM_SIZE = 10
                            # BIGGER_SIZE = 18
                            # plt.rc('font', size=SMALL_SIZE)          # controls default text sizes
                            # plt.rc('axes', titlesize=SMALL_SIZE)     # fontsize of the axes title
                            # plt.rc('axes', labelsize=SMALL_SIZE)    # fontsize of the x and y labels
                            # plt.rc('xtick', labelsize=SMALL_SIZE)    # fontsize of the tick labels
                            # plt.rc('ytick', labelsize=SMALL_SIZE)    # fontsize of the tick labels
                            # plt.rc('legend', fontsize=MEDIUM_SIZE)    # legend fontsize
                            # plt.rc('figure', titlesize=BIGGER_SIZE)  # fontsize of the figure title
                            sns.set(rc={'figure.figsize':(10,8)})
                            plt.tick_params(axis='both', which='major', labelsize=5)
                            st.write(sns.barplot(x=feature_imp, y=feature_imp.index))
                            plt.xlabel('Feature Importance Score', fontsize=14)
                            # plt.ylabel('Features', fontsize=9)
                            plt.title("Visualizing Important Features for {}".format(dimredvar), fontsize=18, pad=15)

                            buf = BytesIO()
                            plt.savefig(buf, format="png")
                            # st.pyplot(fig)
                            st.image(buf)

                        with col3:
                            pass

                        from sklearn.feature_selection import SelectFromModel

                        chosenthreshold = st.number_input("Please select a threshold level of confidence to retain features (we suggest 0.05, or 5%):", value=0.05, min_value=0.0, max_value=1.0, step=0.01)

                        selector = SelectFromModel(rf, threshold=chosenthreshold)
                        features_important = selector.fit_transform(X, y)

                        st.write('Sample of data with initial features:')
                        st.write(pd.DataFrame(X, columns=dimredvars).head())
                        st.write("")
                        st.info(f"Chosen level of significance to retain features: {chosenthreshold}")
                        st.write("")
                        # st.write('Sample of data with selected features:')
                        reduceddf = pd.DataFrame(features_important)
                        st.session_state.reduceddf = reduceddf
                        # st.write(st.session_state.reduceddf.head())

                        rf = RandomForestClassifier(n_estimators=100, max_depth=3,
                                            bootstrap=True, n_jobs=-1,
                                            random_state=0)
                        rf.fit(features_important, y)

                        # for i in feature_imp:
                        #     reduceddf.columns[i] = feature_imp[i]
                        # print(reduceddf.head())

                        for i, j in zip(reduceddf.columns, feature_imp.index):
                            print(i, j)

                        numfeatures = len(reduceddf.columns)

                        featurestokeep = feature_imp.index[0:numfeatures]

                        print(featurestokeep)

                        # importantdf = df[df.columns.intersection(list(feature_imp.index))]
                        # importantdf = df.drop(columns=[col for col in df if col not in list(feature_imp.index)])
                        # print(importantdf.head())

                        importantdf = df
                        # for i in feature_imp.index:
                            # cols = [c for c in df.columns if c != feature_imp.index[i]]
                            # importantdf = df.loc[:,~df.columns.str.contains(feature_imp.index[i])]
                            # importantdf = importantdf.drop(importantdf.filter(i!=j).columns, axis=1)
                            # importantdf = importantdf.drop(columns=[col for col in importantdf if col not in list(feature_imp.index) and col != dimredvar])
                            # importantdf = importantdf.drop(columns=[col for col in importantdf if col != i and col != dimredvar], inplace=True)
                            # importantdf = importantdf.filter(like=i,axis=1)
                            # importantdf=importantdf.drop(importantdf.filter(like=i,axis=1).columns,axis=1)
                        # print(importantdf.head())

                        st.session_state.importantdf = importantdf.drop(columns=[col for col in importantdf if col not in featurestokeep and col != dimredvar])
                        # print(importantdf.head())

                        st.write('Sample of data with only important features retained per the chosen threshold:')
                        st.write(st.session_state.importantdf.head())

                        # Spawn a new Quill editor
                        st.subheader("Notes on dimension reduction analysis")
                        dimredcontent = st_quill(placeholder="Write your notes here", key="dimredquill")

                        st.session_state.dimrednotes = dimredcontent

                        st.write("Dimension reduction took ", time.time() - start_time, "seconds to run")
                
            elif value == "Regression analysis":

                st.subheader("Linear regression analysis")

                with st.spinner('Please wait while we conduct the linear regression analysis'):

                    my_bar = st.progress(0)

                    time.sleep(5)

                    for percent_complete in range(100):
                        time.sleep(0.2)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    ### From Jupyter - 2. Linear regression

                    # Choose predicted variable - this will become dynamic in the app
                    regvars = list(df.columns)
                    # regvars.insert(0, "NONE")
                    # y = df['BRAND'] # old predicted variable
                    chosenvar = st.selectbox("Please select variable for analysis", ["Choose dependent variable"] + regvars)

                    if chosenvar != "Choose dependent variable":

                        y = df[chosenvar]
                        df3 = y.copy()

                        # Define predictor variables
                        df2 = df.drop([chosenvar], axis=1)
                        indepvars = st.multiselect("Please select independent (input) variables for analysis", regvars)
                        # df3 = df2.copy()
                        # df3 = y.copy()

                        if indepvars != []:

                            for indepvar in indepvars:
                                dfi = df[indepvar]
                                df3 = pd.concat((df3, dfi), axis=1)
                            x = df3.drop([chosenvar], axis=1)
                            # print(df3)
                            # x = df3.iloc[:, :]

                        if chosenvar != "NONE" and indepvars != []: 

                            x, y = np.array(x), np.array(y)

                            x = sm.add_constant(x)

                            model = sm.OLS(y, x)

                            results = model.fit()

                            st.write(results.summary())

                            st.write("")

                            st.write('Predicted response:', results.fittedvalues, sep='\n') # Or print('predicted response:', results.predict(x), sep='\n')

                            ### End

                            st.write("")

                            with st.expander("Interpretation guide"):
                                st.write("Regression searches for relationships (mathematical dependencies) among variables. This differs from correlation because regression analysis seeks to unpack **causal** relationships, where a change in one variable causes another variable to change. We find a function that maps some features or variables to others sufficiently well.\n\n"
                                "The dependent features are called the dependent variables, outputs, or responses. The independent features are called the independent variables, inputs, or predictors.\n\n"
                                "Regression problems usually have one continuous and unbounded dependent variable. The inputs, however, can be continuous, discrete, or even categorical data such as gender, nationality, brand, and so on. It is a common practice to denote the outputs with y and inputs with x.\n\n"
                                "Regression is also useful when you want to forecast a response using a new set of predictors. The regression coefficients can be applied as a multiplier factor to understand how changes in an predictor feature would cause the predicted variable to change.")
                                st.info("Regression coefficients indicate the predicted change in the output variable caused by a change in the input variable.")
                                st.info("The estimated or predicted response for each observation should be as close as possible to the corresponding actual response in the underlying data (in real life). The differences for all observations from the actual responses are called the residuals. Regression is about determining the best predicted weights, that is the weights corresponding to the smallest residuals. To get the best weights, you usually minimize the sum of squared residuals (SSR) for all observations. This approach is called the method of ordinary least squares.")
                                st.info("The variation of actual responses occurs partly due to the dependence on the predictors. However, there is also an additional, systematic, inherent variance of the output. The coefficient of determination, denoted as R-squared, tells you which amount of variation in 𝑦 can be explained by the dependence on 𝐱 using the particular regression model. Larger 𝑅² indicates a better fit and means that the model can better explain the variation of the output with different inputs.")
                                st.error("You should, however, be aware of two problems that might follow the choice of the degree: underfitting and overfitting.\n"
                                "Underfitting occurs when a model can't accurately capture the dependencies among data, usually as a consequence of its own simplicity. It often yields a low 𝑅² with known data and bad generalization capabilities when applied with new data. Overfitting happens when a model learns both dependencies among data and random fluctuations. In other words, a model learns the existing data too well. Complex models, which have many features or terms, are often prone to overfitting. When applied to known data, such models usually yield high 𝑅². However, they often don't generalize well and have significantly lower 𝑅² when used with new data.")
                                st.success("The coefficients in our model indicate the extent to which we can expect the predicted variable chosen to change for a 1 unit change in each respective input variable. The p-values indicate the level of statistical significance (lower means more statistically significant). R-squared is a measure of how well the model explains variance in the data.")
                
                    # Spawn a new Quill editor
                    st.subheader("Notes on linear regression analysis")
                    regcontent = st_quill(placeholder="Write your notes here", key="regquill")

                    st.session_state.regnotes = regcontent

                    st.write("Linear regression took ", time.time() - start_time, "seconds to run")

            elif value == "Clustering analysis":

                st.subheader("Cluster analysis")

                st.info("We run cluster analysis for any variable, with optional grouping by brand, for a more granular analysis, if desired.")

                with st.spinner('Please wait while we conduct the cluster analysis using the K-means algorithm'):

                    my_bar = st.progress(0)

                    time.sleep(10)

                    for percent_complete in range(100):
                        time.sleep(0.01)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                with st.expander("Assumptions and operation of the model"):

                    st.write("The technique of cluster analysis aims to group similar observations in a dataset, such that observations in the same group are as similar to each other as possible, and similarly, observations in different groups are as different to each other as possible.")
                    st.warning("This allows us to get a sense of how many different groups with different characteristics and observable behaviour we have in our data, with respect to a chosen variable.")
                    st.info("Cluster analysis does this by minimizing the distance between observations (or the error involved in grouping them), by considering the mean of each potential group. Each group, by construct, has specific characteristics as it relates to the chosen variable.")
                    st.write("Cluster analysis in this case assumes each explanatory variable has the same within-group variance, with spherical variance, and that clusters are roughly similarly sized. Our data is sourced in a way that produces usable variance, and cluster size as reported can be used to further refine results by choosing different cluster sizes as a tuned hyperparameter.")

                ### From Jupyter - 4. Clustering - K-means

                # Perform K-means clustering for consumer segmentation

                # Size of the data set after removing NaN
                print(df2.shape)
                # Explore type of data and feature names
                #print(df2.sample(5))

                # Select explanatory variable

                # Choose brand

                brands = pd.Series(df2['BRAND'].drop_duplicates()).sort_values().tolist() # Faster than .unique()
                brandsmod = ["All"] + brands

                st.markdown("**Please choose brands. If option 'All' is selected, or no option is selected, the analysis will be run for all brands.**")
                st.info("Tip: it is instructive to compare analysis between brands for different variables.")
                # brandssel = []
                brandssel = st.multiselect("Please select brands of interest.", brandsmod, key="brands")

                if "All" not in brandssel:
                    df = df[df['BRANDNAME'].isin(brands)]

                # if brandssel == []:
                #     pass

                # print(df)

                # Select continuous variables for clustering - this is for age
                # X2 = df2.iloc[0:10000, 14:15] #subsamping for efficiency and speed
                clustervars = list(df.columns)
                chosenvar = st.selectbox("Please select variable for cluster analysis", clustervars)
                if chosenvar == "BRANDNAME":
                    chosenvar = "BRAND"
                chosenvarindex = df.columns.get_loc(chosenvar)

                # X2 = df.iloc[:, 10:11].sample(10000) #subsamping for efficiency and speed
                X2 = df.iloc[:, chosenvarindex:chosenvarindex+1].sample(10000, replace=True) # Subsampling for efficiency and speed

                # Find optimal number of clusters

                # 1. Elbow method
                # Calculate distortions
                distortions = []

                for i in range(1, 16):
                    km = KMeans(n_clusters=i, init='k-means++', n_init=10, 
                                max_iter=300,tol=1e-04, random_state=0)
                    km.fit(X2)
                    distortions.append(sum(np.min(cdist(X2, km.cluster_centers_, 
                                    'euclidean'),axis=1)) / X2.shape[0])

                # Plot distortions
                fig, ax = plt.subplots(figsize=(10, 6))
                plt.plot(range(1, 16), distortions, marker='o')
                plt.xlabel('Number of clusters')
                plt.ylabel('Distortion')
                buf = BytesIO()
                fig.savefig(buf, format="png")
                # st.pyplot(fig)
                st.image(buf)

                st.info("You can use the elbow method to refine the analysis, by informing the number of clusters to use below. You should choose the number of clusters based on where the graph sensibly starts to taper/level off, meaning that additional clusters add little additional explanatory power to the model.")

                # 2. Silhouette method
                sil = []
                kmax = 10
                nclusters = []

                for k in range(2, kmax+1):
                    kmeans = KMeans(n_clusters = k).fit(X2)
                    labels = kmeans.labels_
                    silscore = silhouette_score(X2, labels, metric = 'euclidean')
                    sil.append(silscore)
                    if silscore >= max(sil):
                        nclusters = k
                        print(nclusters)

                # Calculate best silhouette score
                print(max(sil))

                # Plot
                fig, ax = plt.subplots(figsize=(10, 6))
                plt.plot(range(2, kmax+1), sil, marker='o')
                plt.xlabel('Number of clusters')
                plt.ylabel('Silhouette score')
                buf = BytesIO()
                fig.savefig(buf, format="png")
                # st.pyplot(fig)
                st.image(buf)
                # Use the output from the elbow or silhouette method to decide how many clusters to use.
                # Cluster the data

                km = KMeans(n_clusters=3, init='k-means++', 
                            n_init=10, max_iter=300, 
                            tol=1e-04, random_state=0)
                km.fit(X2)

                nclustersused = 0

                silrun = st.button("Run cluster analysis with maximum silhouette score")

                if silrun == True:

                    nclustersused = nclusters

                    # Re-cluster with max silhouette score - this will become dynamic in the app - used to be 4 here

                    X2new = X2.copy()
                    kmnew = KMeans(n_clusters=nclusters, init='k-means++', 
                                n_init=10, max_iter=300, 
                                tol=1e-04, random_state=0)
                    kmnew.fit(X2new) 
                    # Check how many observations are in each cluster

                    # print("Cluster 0 size: %s \nCluster 1 size: %s"
                    #       % (len(km.labels_)- km.labels_.sum(), km.labels_.sum()))
                    # Check cluster size once re-clustered

                    print(kmnew.labels_)

                    Xnew2 = X2.copy()
                    Xnew2["CLUSTERS"] = kmnew.labels_
                    Xnew2.sample(8, random_state=0)

                    for i in range(0, nclusters):
                        countclusters = Xnew2.apply(lambda x: True if x['CLUSTERS'] == i else False , axis=1)
                        # Count number of True in series
                        numOfRows = len(countclusters[countclusters == True].index)
                        st.write('Number of Rows in dataframe in which cluster = {} : '.format(i), numOfRows)

                    Xnew = X2.copy()
                    Xnew["CLUSTERS"] = km.labels_
                    Xnew.sample(8, random_state=0)

                    # Xnew2.drop_duplicates(subset=['brand'])
                    Xnew2 = Xnew2[~Xnew2.index.duplicated(keep='first')]
                    Xnew2.to_csv('Xnew2.csv')

                    # Plot the following variables and their clusters
                    var = [chosenvar]

                    # Plot using seaborn

                    # fig, ax = plt.subplots(figsize=(7, 3))
                    fig = sns.pairplot(Xnew2, vars=var, hue="CLUSTERS", palette=sns.color_palette("hls", nclusters), height=5)
                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    #st.pyplot(fig)
                    st.image(buf)
                    #df.head()
                    #df = df.dropna()

                    if nclustersused != 0:

                        st.success("Our cluster analysis indicates that there are " + str(nclustersused) + " clusters in the data for " + chosenvar + " for the brands selected. Each group has specific characteristics and behaviour that be intuited (you can assign meanings to each group based on hypotheses), and further investigation using data exploration, correlation analysis, and linear/logistic regression analysis can provide further insight for " + chosenvar + " clusters in these brands, as compared to other brands or the whole dataset.")

                st.write("Or choose the number of clusters based on the elbow method, enter this in the slider below, and click the 'Conduct analysis' button")
                numclus = st.slider("Please choose number of clusters", min_value=1, max_value=nclusters)
                ownclus = st.button("Conduct analysis")
                if ownclus == True:

                    nclustersused = numclus
                    
                    # Re-cluster with chosen number of clusters

                    X2new = X2.copy()
                    kmnew = KMeans(n_clusters=numclus, init='k-means++', 
                                n_init=10, max_iter=300, 
                                tol=1e-04, random_state=0)
                    kmnew.fit(X2new) 
                    # Check how many observations are in each cluster

                    # print("Cluster 0 size: %s \nCluster 1 size: %s"
                    #       % (len(km.labels_)- km.labels_.sum(), km.labels_.sum()))
                    # Check cluster size once re-clustered

                    print(kmnew.labels_)

                    Xnew2 = X2.copy()
                    Xnew2["CLUSTERS"] = kmnew.labels_
                    Xnew2.sample(8, random_state=0)

                    for i in range(0, numclus):
                        countclusters = Xnew2.apply(lambda x: True if x['CLUSTERS'] == i else False , axis=1)
                        # Count number of True in series
                        numOfRows = len(countclusters[countclusters == True].index)
                        st.write('Number of Rows in dataframe in which cluster = {} : '.format(i), numOfRows)

                    Xnew = X2.copy()
                    Xnew["CLUSTERS"] = km.labels_
                    Xnew.sample(8, random_state=0)

                    Xnew2 = Xnew2[~Xnew2.index.duplicated(keep='first')]
                    Xnew2.to_csv('Xnew2.csv')

                    # Plot the following variables and their clusters
                    var = [chosenvar]

                    # Plot using seaborn

                    # fig, ax = plt.subplots(figsize=(7, 3))
                    fig = sns.pairplot(Xnew2, vars=var, hue="CLUSTERS", palette=sns.color_palette("hls", numclus), height=5)
                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    #st.pyplot(fig)
                    st.image(buf)
                    #df.head()
                    #df = df.dropna()

                    if nclustersused != 0:

                        st.success("You chose " + str(nclustersused) + " clusters in the data for " + chosenvar + " for the brands selected. Each group has specific characteristics and behaviour that be intuited (you can assign meanings to each group based on hypotheses), and further investigation using data exploration, correlation analysis, and linear/logistic regression analysis can provide further insight for " + chosenvar + " clusters in these brands, as compared to other brands or the whole dataset.")

                # Optional - log transformations

                ### End

                st.write("")

                st.subheader("Notes on cluster analysis")

                # Spawn a new Quill editor
                st.subheader("Notes on cluster analysis")
                clustercontent = st_quill(placeholder="Write your notes here", key="clusterquill")

                st.session_state.clusternotes = clustercontent

                st.write("Running the cluster analysis using K-means took ", time.time() - start_time, "seconds to run")

            elif value == "Conjoint analysis":

                st.subheader("Conjoint analysis")

                st.warning("This model takes some time to run. Please be patient.")

                with st.spinner('Please wait while we conduct the conjoint analysis'):

                    my_bar = st.progress(0)

                    time.sleep(10)

                    for percent_complete in range(100):
                        time.sleep(0.01)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    ### From Jupyter - conjoint analysis

                    st.info("Conjoint analysis is traditionally a method for determining what features of a product are most important to consumers. We extend the usage here to enable us to determine what factors are most important in determining a consumer's choice for a particular brand purchase.")

                    brands = pd.Series(branddf['BRAND NAME'].drop_duplicates()).sort_values().tolist() # Faster than .unique()

                    selbrand = st.selectbox("Please select a brand to analyze", brands)

                    # brandcode = branddf.query('BRAND NAME == %s' % str(selbrand))['BRAND']
                    brandcode = branddf.loc[branddf['BRAND NAME'] == selbrand, 'BRAND'].item()

                    y = df2.BRAND.apply(lambda x: brandcode if selbrand in x else 0).head(10000) # Dynamic to brand, subsampling for speed and efficiency
                    x = df2[[x for x in df2.columns if x != 'BRAND']].head(10000) # and x !="FIRST_INTERACTION" and x !="LAST_INTERACTION" and x != "DAYS_LEFT_TO_ENGAGE" and x != "FIRST_NAME" and x != "BIRTH_DATE" and x != "PURCHASE_DESCRIPTION"]].head(10000) # Subsampling for speed and efficiency

                    xdum = pd.get_dummies(x, columns=[c for c in x.columns if c != 'BRAND'])
                    xdum.head()

                    # st.write(y.head(100))

                    plt.style.use('bmh')

                    res = sm.OLS(y.astype(float), xdum.astype(float), family=sm.families.Binomial()).fit()
                    # st.subheader("OLS regression results showing importance of each type of factor")
                    # st.write(res.summary())

                    st.subheader("Factor importances:")

                    df_res = pd.DataFrame({
                    'param_name': res.params.keys()
                    , 'param_w': res.params.values
                    , 'pval': res.pvalues
                    })
                    # adding field for absolute of parameters
                    df_res['abs_param_w'] = np.abs(df_res['param_w'])
                    # marking field is significant under 95% confidence interval
                    df_res['is_sig_95'] = (df_res['pval'] < 0.05)
                    # constructing color naming for each param
                    df_res['c'] = ['blue' if x else 'red' for x in df_res['is_sig_95']]

                    # make it sorted by abs of parameter value
                    df_res = df_res.sort_values(by='abs_param_w', ascending=True)

                    fig, ax = plt.subplots(figsize=(14, 8))
                    plt.title('Fcator importances') # Used to be Part Worth
                    pwu = df_res['param_w']
                    xbar = np.arange(len(pwu))
                    plt.barh(xbar, pwu, color=df_res['c'])
                    plt.yticks(xbar, labels=df_res['param_name'])
                    # plt.show()

                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    st.image(buf)

                    st.subheader("Absolute and relative/normalized importances:")

                    # need to assemble per attribute for every level of that attribute in dicionary
                    range_per_feature = dict()
                    for key, coeff in res.params.items():
                        sk =  key.split('_')
                        feature = sk[0]
                        if len(sk) == 1:
                            feature = key
                        if feature not in range_per_feature:
                            range_per_feature[feature] = list()
                            
                        range_per_feature[feature].append(coeff)
                    # importance per feature is range of coef in a feature
                    # while range is simply max(x) - min(x)
                    importance_per_feature = {
                        k: max(v) - min(v) for k, v in range_per_feature.items()
                    }

                    # compute relative importance per feature
                    # or normalized feature importance by dividing 
                    # sum of importance for all features
                    total_feature_importance = sum(importance_per_feature.values())
                    relative_importance_per_feature = {
                        k: 100 * round(v/total_feature_importance, 3) for k, v in importance_per_feature.items()
                    }

                    alt_data = pd.DataFrame(
                        list(importance_per_feature.items()), 
                        columns=['attr', 'importance']
                    ).sort_values(by='importance', ascending=False)

                    fig, ax = plt.subplots(figsize=(12, 8))
                    xbar = np.arange(len(alt_data['attr']))
                    plt.title('Importance')
                    plt.barh(xbar, alt_data['importance'])
                    for i, v in enumerate(alt_data['importance']):
                        ax.text(v , i + .25, '{:.2f}'.format(v))
                    plt.ylabel('attributes')
                    plt.xlabel('% importance')
                    plt.yticks(xbar, alt_data['attr'])
                    plt.show()

                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    st.image(buf)

                    alt_data = pd.DataFrame(
                        list(relative_importance_per_feature.items()), 
                        columns=['attr', 'relative_importance (pct)']
                    ).sort_values(by='relative_importance (pct)', ascending=False)

                    fig, ax = plt.subplots(figsize=(12, 8))
                    xbar = np.arange(len(alt_data['attr']))
                    plt.title('Relative importance / Normalized importance')
                    plt.barh(xbar, alt_data['relative_importance (pct)'])
                    for i, v in enumerate(alt_data['relative_importance (pct)']):
                        ax.text(v , i + .25, '{:.2f}%'.format(v))
                    plt.ylabel('attributes')
                    plt.xlabel('% relative importance')
                    plt.yticks(xbar, alt_data['attr'])
                    plt.show()

                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    st.image(buf)

                    st.write("")

                    st.success("The results of the conjoint analysis show us which factors should be paid particular attention, in framing a consumer's choice of purchase for the chosen brand, " + selbrand + ".")

                    # Spawn a new Quill editor
                    st.subheader("Notes on conjoint analysis")
                    conjointcontent = st_quill(placeholder="Write your notes here", key="conjointquill")

                    st.session_state.conjointnotes = conjointcontent

                    st.write("Running the conjoint analysis took ", time.time() - start_time, "seconds to run")

            elif value == "Decision trees":

                st.subheader("Decision tree analysis")

                with st.spinner('Please wait while we conduct the decision tree analysis'):

                    my_bar = st.progress(0)

                    time.sleep(10)

                    for percent_complete in range(100):
                        time.sleep(0.01)
                        my_bar.progress(percent_complete + 1)

                    start_time = time.time()

                    st.info("A decision tree is a decision support tool that uses a tree-like model of decisions and their possible consequences, including chance event outcomes, resource costs, and utility. It is one way to display an algorithm that only contains conditional control statements.")
                    st.warning("In this way, a decision tree shows how a particular decision was made, using particular explanatory variables, in our case, how a particular brand was chosen based on gender, age, and city of consumer.")

                    brands = pd.Series(branddf['BRAND NAME'].drop_duplicates()).sort_values().tolist() # Faster than .unique()

                    selbrand = st.selectbox("Please select a brand to analyze", brands)

                    branddf2 = branddf[~branddf["BRAND"].duplicated(keep='first')]

                    # brandcode = branddf2.loc[branddf2['BRAND NAME'] == selbrand, 'BRAND'].item()
                    brandcode = branddf.loc[branddf['BRAND NAME'] == selbrand, 'BRAND'].iloc[0]
                    print(brandcode)

                    # Split data into features (X) and response (y)
                    X = df.loc[0:10000,["GENDER", "AGE", "CITY"]] # Subsampling for speed and efficiency, dynamic to certain relevant explanatory variables (could combine with dimensionality reduction or with conjoint analysis above)
                    y = df2.BRAND.apply(lambda x: brandcode if selbrand in x else 0).head(10001) # Dynamic to brand, subsampling for speed and efficiency

                    # Fit data to tree-based regression model
                    regressor = DecisionTreeRegressor(random_state=0)
                    regressor=regressor.fit(X,y)

                    # Visualising the decision tree regression results
                    plt.figure(figsize=(6,6), dpi=150)
                    plot_tree(regressor,max_depth=3,feature_names=X.columns, impurity=False)
                    plt.show()

                    buf = BytesIO()
                    plt.savefig(buf, format="png")
                    # st.image(buf)

                    # A scatter plot of latitude vs longitude
                    df.plot.scatter(x='AGE',y='GENDER',c='DarkBlue',s=1.5) # Dynamic to chosen variables

                    plt.figure(figsize=[6,4], dpi=120)
                    cutx, cuty = 50, 50
                    plt.ylim(0,cuty)       
                    plt.xlim(0,cutx)
                    plt.xlabel('AGE')
                    plt.ylabel('GENDER')
                    plt.scatter(x=X['AGE'],y=X['GENDER'],c=df['BRAND'].head(10001)) # Dynamic to chosen variables

                    buf2 = BytesIO()
                    plt.savefig(buf2, format="png")
                    st.image(buf2)

                    # splits = regressor.tree_.threshold[:2]
                    # print(splits, cutx, cuty)
                    # plt.plot([splits[1],splits[1]], [0,cuty]) 
                    # plt.plot([splits[1],cutx], [splits[0],splits[0]])
                    # plt.colorbar()
                    # plt.show()

                    # buf3 = BytesIO()
                    # plt.savefig(buf3, format="png")
                    # st.image(buf3)

                    X, y = make_regression(n_samples=1000, n_features=2,n_informative=2,
                                    random_state=0)
                    reg = DecisionTreeRegressor(max_depth=3).fit(X, y)

                    # Parameters
                    plot_colors = "ryb"
                    plot_step = 0.02

                    # Plot the decision boundary
                    f, axes =plt.subplots(ncols=1,nrows=2,figsize=(30, 30))

                    x_min, x_max = X[:, 0].min() - 1, X[:, 0].max() + 1
                    y_min, y_max = X[:, 1].min() - 1, X[:, 1].max() + 1
                    xx, yy = np.meshgrid(np.arange(x_min, x_max, plot_step),
                                        np.arange(y_min, y_max, plot_step))
                    plt.tight_layout(h_pad=0.5, w_pad=0.5, pad=2.5)

                    Z = reg.predict(np.c_[xx.ravel(), yy.ravel()])
                    Z = Z.reshape(xx.shape)
                    cs = plt.contourf(xx, yy, Z, cmap=plt.cm.Blues)

                    plt.xlabel('AGE')
                    plt.ylabel('GENDER')

                    axes[1].scatter(X[:, 0], X[:, 1], c=y,
                                cmap='Oranges', edgecolor='black', s=15)

                    plot_tree(reg, filled=True, feature_names=['AGE', 'GENDER'],
                            ax=axes[0], fontsize=10,
                            class_names='Target')

                    plt.show()

                    buf3 = BytesIO()
                    plt.savefig(buf3, format="png")
                    st.image(buf3)

                    st.write("")

                    st.success("We interpret the decision tree by analyzing how a consumer moves downwards and to the left, meaning that each condition referenced is true if we move to the bottom left, whereas as we move to the bottom right in any particular branch of the tree the relevant referenced condition is false, or opposite to that stated.")

                    # Spawn a new Quill editor
                    st.subheader("Notes on decision tree analysis")
                    dectreecontent = st_quill(placeholder="Write your notes here", key="dectreequill")

                    st.session_state.dectreenotes = dectreecontent

                    st.write("Conducting the decision tree analysis took ", time.time() - start_time, "seconds to run")

            else:

                st.warning('Please select a methodology above')

            with st.expander("Build results presentation", expanded=False):
            
                st.info("When finished, click below to build the presentation with your results")

                buildpres = st.button("Build presentation")

                if buildpres == True:

                    # Building Powerpoint presentation

                    from pptx import Presentation
                    from pptx.enum.shapes import MSO_SHAPE
                    from pptx.dml.color import RGBColor
                    from pptx.util import Inches, Pt
                    from pptx.enum.dml import MSO_THEME_COLOR
                    title='   Analytics Playground\n\
                    Results from analysis'
                    APlogo='./Powerpoint/APlogo.png'
                    ABIlogo='./Powerpoint/ABIlogo.png'
                    prs = Presentation()

                    # Slide 1

                    # Add colour bar

                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    prs.slide_width = Inches(16)
                    prs.slide_height = Inches(9)
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
                    )
                    shape.shadow.inherit = False
                    fill=shape.fill
                    fill.solid()
                    fill.fore_color.rgb=RGBColor(255,0,0)
                    shape.text= title
                    line=shape.line
                    line.color.rgb=RGBColor(255,0,0)
                    logo1=slide.shapes.add_picture(APlogo,Inches(13.5),Inches(6.0),height=Inches(1.08),width=Inches(1.0))
                    logo2=slide.shapes.add_picture(ABIlogo,Inches(14.5),Inches(5.8),height=Inches(1.5),width=Inches(1.51))

                    # Add text box for results

                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
                    shape.shadow.inherit = False
                    fill=shape.fill
                    fill.solid()
                    fill.fore_color.rgb=RGBColor(255,0,0)
                    shape.text= "   Results from Exploratory Data Analysis"
                    line=shape.line
                    line.color.rgb=RGBColor(255,0,0)
                    logo1=slide.shapes.add_picture(APlogo,Inches(14.5),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                    logo2=slide.shapes.add_picture(ABIlogo,Inches(15.0),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(5)
                    height = Inches(5)
                    text_box=slide.shapes.add_textbox(left, top, width, height)
                    tb=text_box.text_frame
                    tb.text = st.session_state.edanotes
                    prg=tb.add_paragraph()
                    prg.text=" "
                    prg=tb.add_paragraph()
                    prg.text=''
                    correlpic = slide.shapes.add_picture('correl.jpg', Inches(8), Inches(1.3), height=Inches(3.7), width=Inches(6.3))
                    ppspic = slide.shapes.add_picture('pps.jpg', Inches(8), Inches(5.1), height=Inches(3.7), width=Inches(7.3))

                    prs.save('EDA_presentation.pptx')

                    os.startfile("EDA_presentation.pptx")

                